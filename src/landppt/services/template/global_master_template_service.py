"""
Global Master Template Service for managing reusable master templates
"""

import json
import logging
import time
import base64
from collections import Counter
from typing import Dict, Any, List, Optional
from io import BytesIO
from sqlalchemy.exc import IntegrityError

from ...ai import get_ai_provider, get_role_provider, AIMessage, MessageRole
from ...ai.base import TextContent, ImageContent, MessageContentType
from ...core.config import ai_config
from ...database.service import DatabaseService
from ...database.database import AsyncSessionLocal
from ..prompts.system_prompts import SystemPrompts
from ..prompts.template_prompts import TemplatePrompts

# Configure logger for this module
logger = logging.getLogger(__name__)


class GlobalMasterTemplateService:
    """Service for managing global master templates"""

    def __init__(
        self,
        provider_name: Optional[str] = None,
        user_id: Optional[int] = None,
        allow_system_template_write: bool = False,
    ):
        self.provider_name = provider_name
        self.user_id = user_id
        self.allow_system_template_write = bool(allow_system_template_write)

    @property
    def ai_provider(self):
        """Dynamically get AI provider to ensure latest config"""
        provider, _ = get_role_provider("template", provider_override=self.provider_name)
        return provider

    def _get_template_role_provider(self):
        """Get provider and settings for template generation role (sync version)"""
        return get_role_provider("template", provider_override=self.provider_name)

    async def _get_template_role_provider_async(self):
        """Get provider and settings for template generation role (async version with landppt support)"""
        logger.info(f"_get_template_role_provider_async called with user_id={self.user_id}")
        
        if self.user_id is not None:
            try:
                from ..db_config_service import get_db_config_service, get_user_ai_provider
                
                config_service = get_db_config_service()
                user_config = await config_service.get_all_config(user_id=self.user_id)
                
                # Get template-specific or default provider and model (use the same mapping as AIConfig)
                role_provider_key, role_model_key = ai_config.MODEL_ROLE_FIELDS.get(
                    "template",
                    ("template_generation_model_provider", "template_generation_model_name"),
                )
                provider_name = self.provider_name or user_config.get(role_provider_key) or user_config.get("default_ai_provider") or "landppt"
                model = user_config.get(role_model_key)
                
                if not model:
                    provider_model_key = f"{provider_name}_model"
                    model = user_config.get(provider_model_key)
                
                logger.info(f"Template provider determined: {provider_name}, model: {model}")
                
                settings = {
                    "role": "template",
                    "provider": provider_name,
                    "model": model
                }
                
                # Use get_user_ai_provider which properly handles landppt system-level config
                provider = await get_user_ai_provider(self.user_id, provider_name)
                logger.info(f"Got template provider successfully: {provider_name}, model: {model}")
                return provider, settings
                
            except Exception as e:
                logger.error(f"Failed to get user template provider async: {e}", exc_info=True)
        else:
            logger.warning("user_id is None, falling back to global config")
        
        # Fall back to global config
        logger.warning("Falling back to global get_role_provider")
        return get_role_provider("template", provider_override=self.provider_name)

    async def _text_completion(self, *, prompt: str, **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        prompt = SystemPrompts.with_text_cache_prefix(prompt)
        return await provider.text_completion(prompt=prompt, **kwargs)

    async def _chat_completion(self, *, messages: List[AIMessage], **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        messages = SystemPrompts.normalize_messages_for_cache(messages)
        return await provider.chat_completion(messages=messages, **kwargs)

    async def _stream_text_completion(self, *, prompt: str, **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        prompt = SystemPrompts.with_text_cache_prefix(prompt)
        if hasattr(provider, 'stream_text_completion'):
            async for chunk in provider.stream_text_completion(prompt=prompt, **kwargs):
                yield chunk
        else:
            response = await provider.text_completion(prompt=prompt, **kwargs)
            yield response.content

    async def _stream_chat_completion(self, *, messages: List[AIMessage], **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        messages = SystemPrompts.normalize_messages_for_cache(messages)
        if hasattr(provider, 'stream_chat_completion'):
            async for chunk in provider.stream_chat_completion(messages=messages, **kwargs):
                yield chunk
        else:
            response = await provider.chat_completion(messages=messages, **kwargs)
            yield response.content

    async def create_template(self, template_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create a new global master template"""
        try:
            template_settings = ai_config.get_model_config_for_role("template", provider_override=self.provider_name)
            # Validate required fields
            required_fields = ['template_name', 'html_template']
            for field in required_fields:
                if not template_data.get(field):
                    raise ValueError(f"Missing required field: {field}")

            # Check if template name already exists
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                existing = await db_service.get_global_master_template_by_name(
                    template_data['template_name'],
                    user_id=self.user_id,
                )
                if existing:
                    raise ValueError(f"Template name '{template_data['template_name']}' already exists")

            # Generate preview image if not provided
            if not template_data.get('preview_image'):
                template_data['preview_image'] = await self._generate_preview_image(template_data['html_template'])

            # Extract style config if not provided
            if not template_data.get('style_config'):
                template_data['style_config'] = self._extract_style_config(template_data['html_template'])

            # Set default values
            template_data.setdefault('description', '')
            template_data.setdefault('tags', [])
            template_data.setdefault('is_default', False)
            template_data.setdefault('is_active', True)
            template_data.setdefault('created_by', 'system')
            if self.user_id is not None:
                template_data.setdefault('user_id', self.user_id)

            # Create template
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                template = await db_service.create_global_master_template(
                    template_data,
                    user_id=self.user_id,
                )

                return {
                    "id": template.id,
                    "user_id": template.user_id,
                    "template_name": template.template_name,
                    "description": template.description,
                    "preview_image": template.preview_image,
                    "tags": template.tags,
                    "is_default": template.is_default,
                    "is_active": template.is_active,
                    "usage_count": template.usage_count,
                    "created_by": template.created_by,
                    "created_at": template.created_at,
                    "updated_at": template.updated_at
                }

        except IntegrityError as e:
            logger.warning(f"Template create integrity constraint violation: {e}")
            raise ValueError(f"Template name '{template_data.get('template_name')}' already exists")
        except Exception as e:
            logger.error(f"Failed to create global master template: {e}")
            raise

    async def get_all_templates(self, active_only: bool = True) -> List[Dict[str, Any]]:
        """Get all global master templates"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                templates = await db_service.get_all_global_master_templates(
                    active_only,
                    user_id=self.user_id,
                )

                return [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

        except Exception as e:
            logger.error(f"Failed to get global master templates: {e}")
            raise

    async def get_all_templates_paginated(
        self,
        active_only: bool = True,
        page: int = 1,
        page_size: int = 6,
        search: Optional[str] = None
    ) -> Dict[str, Any]:
        """Get all global master templates with pagination"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)

                # Calculate offset
                offset = (page - 1) * page_size

                # Get templates with pagination
                templates, total_count = await db_service.get_global_master_templates_paginated(
                    active_only=active_only,
                    offset=offset,
                    limit=page_size,
                    search=search,
                    user_id=self.user_id,
                )

                # Calculate pagination info
                total_pages = (total_count + page_size - 1) // page_size
                has_next = page < total_pages
                has_prev = page > 1

                template_list = [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

                return {
                    "templates": template_list,
                    "pagination": {
                        "current_page": page,
                        "page_size": page_size,
                        "total_count": total_count,
                        "total_pages": total_pages,
                        "has_next": has_next,
                        "has_prev": has_prev
                    }
                }
        except Exception as e:
            logger.error(f"Failed to get paginated templates: {e}")
            raise

    async def get_template_by_id(self, template_id: int) -> Optional[Dict[str, Any]]:
        """Get global master template by ID"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                template = await db_service.get_global_master_template_by_id(
                    template_id,
                    user_id=self.user_id,
                )

                if not template:
                    return None

                return {
                    "id": template.id,
                    "user_id": template.user_id,
                    "template_name": template.template_name,
                    "description": template.description,
                    "html_template": template.html_template,
                    "preview_image": template.preview_image,
                    "style_config": template.style_config,
                    "tags": template.tags,
                    "is_default": template.is_default,
                    "is_active": template.is_active,
                    "usage_count": template.usage_count,
                    "created_by": template.created_by,
                    "created_at": template.created_at,
                    "updated_at": template.updated_at
                }

        except Exception as e:
            logger.error(f"Failed to get global master template {template_id}: {e}")
            raise

    async def update_template(self, template_id: int, update_data: Dict[str, Any]) -> bool:
        """Update a global master template"""
        try:
            # Check if template name conflicts (if being updated)
            if 'template_name' in update_data:
                async with AsyncSessionLocal() as session:
                    db_service = DatabaseService(session)
                    existing = await db_service.get_global_master_template_by_name(
                        update_data['template_name'],
                        user_id=self.user_id,
                    )
                    if existing and existing.id != template_id:
                        raise ValueError(f"Template name '{update_data['template_name']}' already exists")

            # Update preview image if HTML template is updated
            if 'html_template' in update_data and 'preview_image' not in update_data:
                update_data['preview_image'] = await self._generate_preview_image(update_data['html_template'])

            # Update style config if HTML template is updated
            if 'html_template' in update_data and 'style_config' not in update_data:
                update_data['style_config'] = self._extract_style_config(update_data['html_template'])

            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                return await db_service.update_global_master_template(
                    template_id,
                    update_data,
                    user_id=self.user_id,
                    allow_system_write=self.allow_system_template_write,
                )

        except IntegrityError as e:
            logger.warning(f"Template update integrity constraint violation: {e}")
            raise ValueError(f"Template name '{update_data.get('template_name')}' already exists")
        except Exception as e:
            logger.error(f"Failed to update global master template {template_id}: {e}")
            raise

    async def delete_template(self, template_id: int) -> bool:
        """Delete a global master template"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)

                # Check if template exists
                template = await db_service.get_global_master_template_by_id(
                    template_id,
                    user_id=self.user_id,
                )
                if not template:
                    logger.warning(f"Template {template_id} not found for deletion")
                    return False

                # Check if it's the default template
                if template.is_default:
                    raise ValueError("Cannot delete the default template")

                logger.info(f"Deleting template {template_id}: {template.template_name}")
                result = await db_service.delete_global_master_template(
                    template_id,
                    user_id=self.user_id,
                    allow_system_write=self.allow_system_template_write,
                )

                if result:
                    logger.info(f"Successfully deleted template {template_id}")
                else:
                    logger.warning(f"Failed to delete template {template_id} - no rows affected")

                return result

        except Exception as e:
            logger.error(f"Failed to delete global master template {template_id}: {e}")
            raise

    async def set_default_template(self, template_id: int) -> bool:
        """Set a template as default"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                return await db_service.set_default_global_master_template(
                    template_id,
                    user_id=self.user_id,
                    allow_system_write=self.allow_system_template_write,
                )

        except Exception as e:
            logger.error(f"Failed to set default template {template_id}: {e}")
            raise

    async def get_default_template(self) -> Optional[Dict[str, Any]]:
        """Get the default template"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                template = await db_service.get_default_global_master_template(user_id=self.user_id)

                if not template:
                    return None

                return {
                    "id": template.id,
                    "user_id": template.user_id,
                    "template_name": template.template_name,
                    "description": template.description,
                    "html_template": template.html_template,
                    "preview_image": template.preview_image,
                    "style_config": template.style_config,
                    "tags": template.tags,
                    "is_default": template.is_default,
                    "is_active": template.is_active,
                    "usage_count": template.usage_count,
                    "created_by": template.created_by,
                    "created_at": template.created_at,
                    "updated_at": template.updated_at
                }

        except Exception as e:
            logger.error(f"Failed to get default template: {e}")
            raise

    @staticmethod
    def _get_template_resource_performance_prompt_text() -> str:
        """统一模板生成阶段的资源可达性与性能约束。"""
        return TemplatePrompts.get_template_resource_performance_prompt_text()

    @staticmethod
    def _get_template_annotation_prompt_text() -> str:
        """固定画布与母版职责分层提示。"""
        return TemplatePrompts.get_template_annotation_prompt_text()

    @staticmethod
    def _get_template_generation_creative_prompt_text() -> str:
        """母版创意目标，强调视觉语言系统而非单页样张。"""
        return TemplatePrompts.get_template_generation_creative_prompt_text()

    @staticmethod
    def _get_template_generation_method_prompt_text() -> str:
        """模板生成的创意思考顺序。"""
        return TemplatePrompts.get_template_generation_method_prompt_text()

    def _get_template_generation_requirements_prompt_text(self) -> str:
        """母版生成技术要求。"""
        return TemplatePrompts.get_template_generation_requirements_prompt_text()

    def _build_template_generation_prompt(self, user_prompt: str, mode_instruction: str = "") -> str:
        """组装模板生成提示词。"""
        return TemplatePrompts.build_template_generation_prompt(user_prompt, mode_instruction=mode_instruction)

    async def generate_template_with_ai(self, prompt: str, template_name: str, description: str = "",
                                      tags: List[str] = None, generation_mode: str = "text_only",
                                      reference_image: dict = None, reference_pptx: dict = None,
                                      prompt_is_ready: bool = False):
        """Generate a new template using AI (non-streaming) - does not save to database"""
        import json

        if generation_mode == "pptx_extract":
            if not reference_pptx:
                raise ValueError("PPTX提取模式需要上传PPTX文件")

            pptx_context = self._extract_pptx_template_reference(reference_pptx)
            extracted_summary = pptx_context.get("analysis_summary", "")
            extracted_image = pptx_context.get("reference_image")

            prompt = (
                f"{prompt}\n\n"
                "请基于以下从上传PPTX中提取的模板信息生成HTML母版模板。"
                "重点提取视觉风格、版式结构、字体与配色规律，不要照搬原始文案内容。\n\n"
                "如果从多页中推断出稳定的母版元素（如页眉、页脚、页码区域），请在生成结果中保留它们的相对位置和风格。\n\n"
                f"{extracted_summary}"
            )

            if extracted_image:
                reference_image = extracted_image
                generation_mode = "reference_style"
            else:
                generation_mode = "text_only"

        # 构建AI提示词
        if generation_mode == "text_only" or not reference_image:
            # 纯文本生成模式
            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(prompt)
            messages = [{"role": "user", "content": ai_prompt}]
        else:
            # 多模态生成模式
            if generation_mode == "reference_style":
                mode_instruction = """
请参考上传图片的气质、配色和版式逻辑，但按 PPT 母版需求重新组织，不要机械复刻单一页面。
"""
            else:  # one_to_one
                mode_instruction = """
请尽量贴近上传图片的风格和版式特征，同时保留可复用的标题、内容和页脚结构。
"""

            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(
                prompt,
                mode_instruction=mode_instruction,
            )

            # 构建多模态消息
            # 确保图片URL格式正确
            image_data = reference_image['data']
            if not image_data.startswith("data:"):
                # 如果是纯base64数据,添加data URL前缀
                image_data = f"data:{reference_image['type']};base64,{image_data}"

            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": ai_prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": image_data
                            }
                        }
                    ]
                }
            ]

        try:
            # 获取模板生成任务的配置（使用异步版本以支持 landppt 系统级配置）
            provider, template_settings = await self._get_template_role_provider_async()

            if not provider:
                raise ValueError("AI服务未配置或不可用")

            # 转换消息格式
            ai_messages = []
            for msg in messages:
                if isinstance(msg["content"], str):
                    # 纯文本消息
                    ai_messages.append(AIMessage(
                        role=MessageRole.USER,
                        content=[TextContent(text=msg["content"])]
                    ))
                else:
                    # 多模态消息
                    content_parts = []
                    for part in msg["content"]:
                        if part["type"] == "text":
                            content_parts.append(TextContent(text=part["text"]))
                        elif part["type"] == "image_url":
                            # 提取图片URL (已经是完整的data URL格式)
                            image_url = part["image_url"]["url"]
                            if image_url.startswith("data:"):
                                content_parts.append(ImageContent(
                                    image_url={"url": image_url},
                                    content_type=MessageContentType.IMAGE_URL
                                ))
                    ai_messages.append(AIMessage(
                        role=MessageRole.USER,
                        content=content_parts
                    ))

            ai_response = await self._chat_completion(
                messages=ai_messages,
                model=template_settings.get('model')
            )
            full_response = ai_response.content

            if not full_response or not full_response.strip():
                raise ValueError("AI服务返回空响应")

            html_template = self._extract_html_from_response(full_response)
            if not html_template or not html_template.strip():
                raise ValueError("AI响应中未找到有效的HTML模板")

            logger.info(f"Generated HTML template length: {len(html_template)}")

            # 返回结果（不保存到数据库）
            return {
                'html_template': html_template,
                'template_name': template_name,
                'description': description or f"AI生成的模板：{prompt[:100]}",
                'tags': tags or ['AI生成'],
                'llm_response': full_response  # 包含完整的LLM响应
            }

        except Exception as e:
            logger.error(f"Failed to generate template with AI: {e}", exc_info=True)
            raise

    async def generate_template_with_ai_stream(self, prompt: str, template_name: str, description: str = "",
                                             tags: List[str] = None, generation_mode: str = "text_only",
                                             reference_image: dict = None, prompt_is_ready: bool = False):
        """Generate a new template using AI with streaming response"""
        import asyncio
        import json

        # 构建AI提示词
        if generation_mode == "text_only" or not reference_image:
            # 纯文本生成模式
            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(prompt)
        else:
            # 多模态生成模式
            if generation_mode == "reference_style":
                mode_instruction = """
请参考图片的气质、配色和版式逻辑，但按 PPT 母版需求重新组织，不要机械复刻单一页面。
"""
            else:  # exact_replica
                mode_instruction = """
请尽量贴近参考图片的风格和版式特征，同时保留可复用的标题、内容和页脚结构。
"""

            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(
                prompt,
                mode_instruction=mode_instruction,
            )

        try:
            # 获取模板生成任务的配置（使用异步版本以支持 landppt 系统级配置）
            provider, template_settings = await self._get_template_role_provider_async()

            # 构建AI消息
            if generation_mode != "text_only" and reference_image:
                # 多模态消息
                # 确保图片URL格式正确 (OpenAI需要完整的data URL格式)
                image_url = reference_image["data"]
                if not image_url.startswith("data:"):
                    # 如果是纯base64数据,添加data URL前缀
                    image_type = reference_image.get("type", "image/png")
                    image_url = f"data:{image_type};base64,{image_url}"

                content_parts = [
                    TextContent(text=ai_prompt),
                    ImageContent(image_url={"url": image_url})
                ]
                messages = [AIMessage(role=MessageRole.USER, content=content_parts)]

                # 检查AI提供商是否支持流式聊天
                if hasattr(provider, 'stream_chat_completion'):
                    # 使用流式聊天API
                    full_response = ""
                    async for chunk in provider.stream_chat_completion(
                        messages=messages,
                        temperature=0.7,
                        model=template_settings.get('model')
                    ):
                        full_response += chunk
                        yield {
                            'type': 'thinking',
                            'content': chunk
                        }
                else:
                    # 使用标准聊天API
                    response = await self._chat_completion(
                        messages=messages,
                        temperature=0.7
                    )
                    full_response = response.content

                    # 模拟流式输出
                    yield {'type': 'thinking', 'content': '🖼️ 正在分析参考图片...\n\n'}
                    await asyncio.sleep(1)
                    yield {'type': 'thinking', 'content': full_response}
            else:
                # 纯文本消息
                if hasattr(provider, 'stream_text_completion'):
                    # 使用流式API
                    full_response = ""
                    async for chunk in provider.stream_text_completion(
                        prompt=ai_prompt,
                        temperature=0.7,
                        model=template_settings.get('model')
                    ):
                        full_response += chunk
                        yield {
                            'type': 'thinking',
                            'content': chunk
                        }
                else:
                    # 使用标准文本完成API
                    response = await provider.text_completion(
                        prompt=ai_prompt,
                        temperature=0.7,
                        model=template_settings.get('model')
                    )
                    full_response = response.content

                    # 模拟流式输出
                    yield {'type': 'thinking', 'content': '🤔 正在分析您的需求...\n\n'}
                    await asyncio.sleep(1)
                    yield {'type': 'thinking', 'content': full_response}

                # 流式完成后，处理完整响应
                yield {'type': 'thinking', 'content': '\n\n✨ 优化样式和交互效果...\n'}
                await asyncio.sleep(0.5)

                # 处理AI响应
                html_template = self._extract_html_from_response(full_response)

                if not html_template or not html_template.strip():
                    raise ValueError("Generated HTML template is empty")

                yield {'type': 'thinking', 'content': '✅ 模板生成完成，准备预览...\n'}
                await asyncio.sleep(0.3)

                # 返回生成完成的信息，包含HTML模板用于预览
                yield {
                    'type': 'complete',
                    'message': '模板生成完成！',
                    'html_template': html_template,
                    'template_name': template_name,
                    'description': description or f"AI生成的模板：{prompt[:100]}",
                    'tags': tags or ['AI生成'],
                    'llm_response': full_response  # 添加完整的LLM响应
                }

        except Exception as e:
            logger.error(f"Failed to generate template with AI stream: {e}", exc_info=True)
            yield {
                'type': 'error',
                'message': str(e)
            }

    async def adjust_template_with_ai_stream(self, current_html: str, adjustment_request: str, template_name: str = "模板"):
        """根据用户反馈调整现有模板，允许做同源升级而不是表面修补。"""
        import asyncio

        ai_prompt = f"""
作为专业的 PPT 模板设计师，请根据用户的调整需求升级现有 HTML 模板。

当前模板：
```html
{current_html}
```

用户调整需求：{adjustment_request}

请按以下要求调整：
- 保留原模板的稳定锚点、占位符、视觉母语和可复用骨架。
- 如果用户需求涉及主舞台表达、内容承载能力、节奏、材质或构图气质，允许对主舞台区进行同源重设计，不要只做表面微调。
- 标题锚点和编号锚点可以做同源优化，但不要改成另一套完全不同的系统；编号锚点不默认在 footer，可位于模板已有的任何稳定位置。
- 如果用户明确要求改变风格方向，可以保留职责分层和占位符，同时重建色彩、字体、材质和组件语气。
- 不要把修改做成局部补丁拼接；调整后的结果仍应像一套完整的新版本母板，而不是在旧页面上缝补。
- 保留所有占位符与可复用结构。
- 输出完整 HTML，继续满足 1280x720、无滚动条，并在固定画布内稳定成立。
{self._get_template_resource_performance_prompt_text()}
{self._get_template_generation_creative_prompt_text()}
{self._get_template_annotation_prompt_text()}

请直接输出调整后的完整 HTML，使用```html```代码块返回。
"""

        try:
            provider, template_settings = await self._get_template_role_provider_async()

            if hasattr(provider, 'stream_text_completion'):
                full_response = ""
                async for chunk in provider.stream_text_completion(
                    prompt=ai_prompt,
                    temperature=0.7,
                    model=template_settings.get('model')
                ):
                    full_response += chunk
                    yield {
                        'type': 'thinking',
                        'content': chunk
                    }

                yield {'type': 'thinking', 'content': '\n\n完成模板调整...\n'}
                await asyncio.sleep(0.5)

                html_template = self._extract_html_from_response(full_response)

                if not html_template or not html_template.strip():
                    raise ValueError("Adjusted HTML template is empty")

                yield {
                    'type': 'complete',
                    'message': '模板调整完成',
                    'html_template': html_template,
                    'template_name': template_name
                }

            else:
                yield {'type': 'thinking', 'content': '正在分析调整需求...\n\n'}
                await asyncio.sleep(1)

                yield {'type': 'thinking', 'content': f'调整需求：{adjustment_request}\n\n'}
                await asyncio.sleep(0.5)

                yield {'type': 'thinking', 'content': '开始重组模板语言...\n'}
                await asyncio.sleep(1)

                response = await provider.text_completion(
                    prompt=ai_prompt,
                    temperature=0.7,
                    model=template_settings.get('model')
                )

                yield {'type': 'thinking', 'content': '完成模板调整...\n'}
                await asyncio.sleep(0.5)

                html_template = self._extract_html_from_response(response.content)

                if not html_template or not html_template.strip():
                    raise ValueError("Adjusted HTML template is empty")

                yield {
                    'type': 'complete',
                    'message': '模板调整完成',
                    'html_template': html_template,
                    'template_name': template_name
                }

        except Exception as e:
            logger.error(f"Failed to adjust template with AI stream: {e}", exc_info=True)
            yield {
                'type': 'error',
                'message': str(e)
            }

    def _decode_uploaded_base64_file(self, raw_data: str) -> bytes:
        """Decode uploaded base64 data (supports raw base64 or data URL)."""
        if not raw_data or not isinstance(raw_data, str):
            raise ValueError("上传文件数据为空")

        data_str = raw_data.strip()
        if data_str.startswith("data:"):
            comma_index = data_str.find(",")
            if comma_index < 0:
                raise ValueError("上传文件数据格式无效")
            data_str = data_str[comma_index + 1 :]

        try:
            return base64.b64decode(data_str, validate=False)
        except Exception as e:
            raise ValueError(f"上传文件Base64解码失败: {e}") from e

    def _safe_pptx_rgb_hex(self, color_obj) -> Optional[str]:
        """Best-effort conversion of python-pptx color object to hex string."""
        try:
            if color_obj is None:
                return None
            rgb = getattr(color_obj, "rgb", None)
            if rgb is None:
                return None
            rgb_text = str(rgb).strip().replace("#", "").replace("0x", "").upper()
            if len(rgb_text) == 6 and all(ch in "0123456789ABCDEF" for ch in rgb_text):
                return f"#{rgb_text}"
        except Exception:
            return None
        return None

    def _is_page_number_like_text(self, text: str) -> bool:
        """Heuristic for page number/footer numbering text in PPTX."""
        import re

        if not text:
            return False
        normalized = str(text).strip()
        if not normalized:
            return False

        normalized = normalized.replace("／", "/").replace("丨", "/").replace("|", "/")
        normalized = re.sub(r"\s+", "", normalized)

        patterns = [
            r"^\d+$",
            r"^\d+/\d+$",
            r"^第?\d+页$",
            r"^第?\d+/\d+页$",
            r"^p\.?\d+$",
            r"^page\d+$",
            r"^\d+-\d+$",
        ]
        lower_text = normalized.lower()
        return any(re.fullmatch(pattern, lower_text) for pattern in patterns)

    def _summarize_common_master_candidates(
        self,
        candidates: List[Dict[str, Any]],
        sampled_slide_count: int,
    ) -> Dict[str, Any]:
        """Summarize stable repeated header/footer/page-number-like elements across sampled slides."""
        if sampled_slide_count <= 1 or not candidates:
            return {"summary_lines": [], "stable_elements": []}

        required_count = max(2, int(sampled_slide_count * 0.6 + 0.999))
        aggregate: Dict[str, Dict[str, Any]] = {}

        for item in candidates:
            signature = str(item.get("signature") or "")
            if not signature:
                continue
            slide_idx = int(item.get("slide_idx") or 0)
            if slide_idx <= 0:
                continue

            slot = aggregate.get(signature)
            if slot is None:
                slot = {
                    "signature": signature,
                    "zone": item.get("zone") or "footer",
                    "kind": item.get("kind") or "shape",
                    "position_desc": item.get("position_desc") or "",
                    "style_hint": item.get("style_hint") or "",
                    "text_examples": [],
                    "slides": set(),
                }
                aggregate[signature] = slot

            if slide_idx in slot["slides"]:
                continue
            slot["slides"].add(slide_idx)

            text_example = (item.get("text_example") or "").strip()
            if text_example and text_example not in slot["text_examples"]:
                slot["text_examples"].append(text_example)

        zone_order = {"header": 0, "footer": 1, "page_number": 2}
        zone_labels = {
            "header": "页眉区域",
            "footer": "页脚区域",
            "page_number": "页码区域",
        }
        kind_labels = {
            "text": "文本框",
            "picture": "图片",
            "shape": "图形",
        }

        stable_elements: List[Dict[str, Any]] = []
        for _, item in aggregate.items():
            count = len(item["slides"])
            if count < required_count:
                continue
            stable_elements.append(
                {
                    "zone": item["zone"],
                    "kind": item["kind"],
                    "position_desc": item["position_desc"],
                    "style_hint": item["style_hint"],
                    "text_examples": item["text_examples"][:3],
                    "slide_hits": count,
                    "sampled_slide_count": sampled_slide_count,
                }
            )

        stable_elements.sort(
            key=lambda x: (
                zone_order.get(str(x.get("zone")), 99),
                -int(x.get("slide_hits") or 0),
                str(x.get("kind") or ""),
                str(x.get("position_desc") or ""),
            )
        )

        if not stable_elements:
            return {"summary_lines": [], "stable_elements": []}

        summary_lines = [
            "- 推断公共母版元素（跨采样页稳定出现，建议在HTML模板中保留相对位置与样式）:"
        ]
        for item in stable_elements[:8]:
            zone_label = zone_labels.get(str(item.get("zone")), str(item.get("zone") or "未知区域"))
            kind_label = kind_labels.get(str(item.get("kind")), str(item.get("kind") or "元素"))
            count = int(item.get("slide_hits") or 0)
            position_desc = str(item.get("position_desc") or "")
            style_hint = str(item.get("style_hint") or "").strip()
            text_examples = item.get("text_examples") or []

            line = f"  - {zone_label}：{kind_label} {position_desc}（{count}/{sampled_slide_count}页）"
            if style_hint:
                line += f"；样式={style_hint}"
            if text_examples and str(item.get('zone')) != "page_number":
                line += f"；示例文本={text_examples[0][:40]}"
            if str(item.get("zone")) == "page_number":
                line += "；疑似页码位置"
            summary_lines.append(line)

        return {"summary_lines": summary_lines, "stable_elements": stable_elements}

    def _extract_pptx_template_reference(self, reference_pptx: Dict[str, Any]) -> Dict[str, Any]:
        """Extract reusable style/layout hints from an uploaded PPTX for AI template generation."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except Exception as e:
            raise ValueError(f"当前环境未安装 python-pptx，无法解析PPTX: {e}") from e

        pptx_bytes = self._decode_uploaded_base64_file(reference_pptx.get("data", ""))
        if not pptx_bytes:
            raise ValueError("PPTX文件为空")

        if len(pptx_bytes) > 50 * 1024 * 1024:
            raise ValueError("PPTX文件过大，请控制在 50MB 以内")

        try:
            prs = Presentation(BytesIO(pptx_bytes))
        except Exception as e:
            raise ValueError(f"PPTX解析失败，请确认文件有效: {e}") from e

        slide_count = len(prs.slides)
        if slide_count == 0:
            raise ValueError("PPTX中没有幻灯片，无法提取模板")

        slide_width = int(getattr(prs, "slide_width", 0) or 0)
        slide_height = int(getattr(prs, "slide_height", 0) or 0)
        emu_per_inch = 914400
        slide_width_px = int(round(slide_width * 96 / emu_per_inch)) if slide_width else 0
        slide_height_px = int(round(slide_height * 96 / emu_per_inch)) if slide_height else 0

        sampled_slides = list(prs.slides)[: min(slide_count, 5)]

        font_counter: Counter = Counter()
        color_counter: Counter = Counter()
        font_size_counter: Counter = Counter()
        layout_counter: Counter = Counter()
        slide_summaries: List[str] = []
        common_master_candidates: List[Dict[str, Any]] = []

        best_picture = None
        best_picture_area = 0

        for slide_idx, slide in enumerate(sampled_slides, start=1):
            shape_count = len(slide.shapes)
            text_box_count = 0
            picture_count = 0
            vector_shape_count = 0
            slide_title = ""
            slide_bg_color = None
            shape_lines: List[str] = []
            seen_master_signatures: set = set()

            try:
                fill = slide.background.fill
                slide_bg_color = self._safe_pptx_rgb_hex(getattr(fill, "fore_color", None))
            except Exception:
                slide_bg_color = None

            for shape in slide.shapes:
                try:
                    shape_type = getattr(shape, "shape_type", None)
                    left = int(getattr(shape, "left", 0) or 0)
                    top = int(getattr(shape, "top", 0) or 0)
                    width = int(getattr(shape, "width", 0) or 0)
                    height = int(getattr(shape, "height", 0) or 0)
                    bbox_desc = ""
                    if slide_width and slide_height:
                        bbox_desc = (
                            f"x={left / slide_width:.0%},y={top / slide_height:.0%},"
                            f"w={width / slide_width:.0%},h={height / slide_height:.0%}"
                        )
                    left_ratio = (left / slide_width) if slide_width else None
                    top_ratio = (top / slide_height) if slide_height else None
                    width_ratio = (width / slide_width) if slide_width else None
                    height_ratio = (height / slide_height) if slide_height else None
                    right_ratio = (left_ratio + width_ratio) if left_ratio is not None and width_ratio is not None else None
                    bottom_ratio = (top_ratio + height_ratio) if top_ratio is not None and height_ratio is not None else None

                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_count += 1
                        area = max(0, width) * max(0, height)
                        try:
                            image = shape.image
                            content_type = getattr(image, "content_type", "") or ""
                            blob = getattr(image, "blob", None)
                            if (
                                blob
                                and isinstance(blob, (bytes, bytearray))
                                and content_type.startswith("image/")
                                and area > best_picture_area
                            ):
                                best_picture_area = area
                                best_picture = {
                                    "blob": bytes(blob),
                                    "content_type": content_type,
                                    "filename": getattr(image, "filename", None) or f"slide_{slide_idx}_image",
                                    "area": area,
                                }
                        except Exception:
                            pass

                        if len(shape_lines) < 8:
                            shape_lines.append(f"- 图片：{bbox_desc}".strip("："))
                        if (
                            top_ratio is not None
                            and bottom_ratio is not None
                            and left_ratio is not None
                            and width_ratio is not None
                            and height_ratio is not None
                        ):
                            zone = None
                            if top_ratio <= 0.22 and bottom_ratio <= 0.35:
                                zone = "header"
                            elif bottom_ratio >= 0.78 and top_ratio >= 0.58:
                                zone = "footer"

                            if zone:
                                sig = (
                                    f"{zone}|picture|"
                                    f"x{round(left_ratio / 0.04) * 0.04:.2f}|"
                                    f"y{round(top_ratio / 0.04) * 0.04:.2f}|"
                                    f"w{round(width_ratio / 0.04) * 0.04:.2f}|"
                                    f"h{round(height_ratio / 0.04) * 0.04:.2f}"
                                )
                                if sig not in seen_master_signatures:
                                    seen_master_signatures.add(sig)
                                    common_master_candidates.append({
                                        "slide_idx": slide_idx,
                                        "zone": zone,
                                        "kind": "picture",
                                        "signature": sig,
                                        "position_desc": bbox_desc or "top/bottom-region",
                                        "style_hint": "",
                                        "text_example": "",
                                    })
                        continue

                    if getattr(shape, "has_text_frame", False):
                        text_box_count += 1
                        text = ""
                        try:
                            text = (shape.text or "").strip()
                        except Exception:
                            text = ""

                        if text and not slide_title:
                            try:
                                is_ph = bool(getattr(shape, "is_placeholder", False))
                            except Exception:
                                is_ph = False
                            if is_ph or top < slide_height * 0.25:
                                slide_title = text.splitlines()[0][:80]

                        text_preview = text.replace("\r", " ").replace("\n", " ").strip()[:80] if text else ""
                        text_meta_parts: List[str] = []
                        primary_font_name = None
                        primary_font_size = None
                        primary_font_color = None
                        try:
                            for paragraph in list(shape.text_frame.paragraphs)[:3]:
                                for run in list(paragraph.runs)[:5]:
                                    font = getattr(run, "font", None)
                                    if not font:
                                        continue
                                    font_name = (getattr(font, "name", None) or "").strip()
                                    if font_name:
                                        font_counter[font_name] += 1
                                        if primary_font_name is None:
                                            primary_font_name = font_name
                                    font_size = getattr(font, "size", None)
                                    if font_size is not None:
                                        try:
                                            pt_value = round(float(font_size.pt), 1)
                                            font_size_counter[f"{pt_value}pt"] += 1
                                            if primary_font_size is None:
                                                primary_font_size = pt_value
                                        except Exception:
                                            pass
                                    font_color = self._safe_pptx_rgb_hex(getattr(font, "color", None))
                                    if font_color:
                                        color_counter[font_color] += 1
                                        if primary_font_color is None:
                                            primary_font_color = font_color
                        except Exception:
                            pass

                        if len(shape_lines) < 8:
                            line = f"- 文本框：{bbox_desc}"
                            if text_preview:
                                line += f" 文本示例=“{text_preview}”"
                            shape_lines.append(line)
                        if (
                            top_ratio is not None
                            and bottom_ratio is not None
                            and left_ratio is not None
                            and width_ratio is not None
                            and height_ratio is not None
                        ):
                            is_top_zone = top_ratio <= 0.22 and bottom_ratio <= 0.35
                            is_bottom_zone = bottom_ratio >= 0.78 and top_ratio >= 0.58
                            is_page_num = (
                                is_bottom_zone
                                and width_ratio <= 0.28
                                and height_ratio <= 0.14
                                and self._is_page_number_like_text(text_preview or text)
                            )

                            zone = None
                            if is_page_num:
                                zone = "page_number"
                            elif is_top_zone:
                                zone = "header"
                            elif is_bottom_zone:
                                zone = "footer"

                            if zone:
                                style_parts = []
                                if primary_font_name:
                                    style_parts.append(f"font={primary_font_name}")
                                if primary_font_size is not None:
                                    style_parts.append(f"size={primary_font_size}pt")
                                if primary_font_color:
                                    style_parts.append(f"color={primary_font_color}")
                                style_hint = ", ".join(style_parts)

                                sig = (
                                    f"{zone}|text|"
                                    f"x{round(left_ratio / 0.04) * 0.04:.2f}|"
                                    f"y{round(top_ratio / 0.04) * 0.04:.2f}|"
                                    f"w{round(width_ratio / 0.04) * 0.04:.2f}|"
                                    f"h{round(height_ratio / 0.04) * 0.04:.2f}|"
                                    f"fs{round((primary_font_size or 0) / 2) * 2 if primary_font_size else 0}|"
                                    f"fn{(primary_font_name or '').lower()[:24]}"
                                )
                                if sig not in seen_master_signatures:
                                    seen_master_signatures.add(sig)
                                    common_master_candidates.append({
                                        "slide_idx": slide_idx,
                                        "zone": zone,
                                        "kind": "text",
                                        "signature": sig,
                                        "position_desc": bbox_desc or "top/bottom-region",
                                        "style_hint": style_hint,
                                        "text_example": "" if is_page_num else (text_preview or ""),
                                    })
                        continue

                    vector_shape_count += 1
                    fill_color = None
                    try:
                        fill = getattr(shape, "fill", None)
                        fill_color = self._safe_pptx_rgb_hex(getattr(fill, "fore_color", None)) if fill else None
                        if fill_color:
                            color_counter[fill_color] += 1
                    except Exception:
                        pass

                    if len(shape_lines) < 8:
                        shape_type_name = str(shape_type).split(".")[-1] if shape_type is not None else "UNKNOWN"
                        shape_lines.append(f"- 图形({shape_type_name})：{bbox_desc}".strip("："))
                    if (
                        top_ratio is not None
                        and bottom_ratio is not None
                        and left_ratio is not None
                        and width_ratio is not None
                        and height_ratio is not None
                    ):
                        zone = None
                        if top_ratio <= 0.22 and bottom_ratio <= 0.35:
                            zone = "header"
                        elif bottom_ratio >= 0.78 and top_ratio >= 0.58:
                            zone = "footer"

                        if zone:
                            sig = (
                                f"{zone}|shape|"
                                f"x{round(left_ratio / 0.04) * 0.04:.2f}|"
                                f"y{round(top_ratio / 0.04) * 0.04:.2f}|"
                                f"w{round(width_ratio / 0.04) * 0.04:.2f}|"
                                f"h{round(height_ratio / 0.04) * 0.04:.2f}|"
                                f"c{(fill_color or '')[:7]}"
                            )
                            if sig not in seen_master_signatures:
                                seen_master_signatures.add(sig)
                                common_master_candidates.append({
                                    "slide_idx": slide_idx,
                                    "zone": zone,
                                    "kind": "shape",
                                    "signature": sig,
                                    "position_desc": bbox_desc or "top/bottom-region",
                                    "style_hint": f"fill={fill_color}" if fill_color else "",
                                    "text_example": "",
                                })
                except Exception:
                    continue

            layout_counter[f"text={text_box_count},pic={picture_count},shape={vector_shape_count}"] += 1

            slide_summary = [
                f"第{slide_idx}页：",
                f"元素数量={shape_count}（文本框{text_box_count}、图片{picture_count}、图形{vector_shape_count}）",
            ]
            if slide_title:
                slide_summary.append(f"标题候选={slide_title}")
            if slide_bg_color:
                slide_summary.append(f"背景色={slide_bg_color}")
            if shape_lines:
                slide_summary.append("布局片段=" + "；".join(shape_lines[:6]))
            slide_summaries.append(" | ".join(slide_summary))

        dominant_fonts = [name for name, _ in font_counter.most_common(5)]
        dominant_colors = [color for color, _ in color_counter.most_common(8)]
        dominant_font_sizes = [size for size, _ in font_size_counter.most_common(6)]
        dominant_layouts = [layout for layout, _ in layout_counter.most_common(3)]
        common_master_summary = self._summarize_common_master_candidates(
            common_master_candidates,
            len(sampled_slides),
        )

        reference_image = None
        if best_picture and best_picture.get("blob") and len(best_picture["blob"]) <= 10 * 1024 * 1024:
            img_b64 = base64.b64encode(best_picture["blob"]).decode("utf-8")
            content_type = best_picture.get("content_type") or "image/png"
            reference_image = {
                "filename": best_picture.get("filename") or "pptx_reference_image",
                "size": len(best_picture["blob"]),
                "type": content_type,
                "data": f"data:{content_type};base64,{img_b64}",
            }

        summary_parts: List[str] = [
            "【PPTX模板提取结果】",
            f"- 文件名：{reference_pptx.get('filename') or 'uploaded.pptx'}",
            f"- 幻灯片总数：{slide_count}",
        ]
        if slide_width_px and slide_height_px:
            summary_parts.append(f"- 页面尺寸（近似像素）：{slide_width_px}x{slide_height_px}")
        if dominant_layouts:
            summary_parts.append(f"- 常见布局结构：{'；'.join(dominant_layouts)}")
        if dominant_fonts:
            summary_parts.append(f"- 高频字体：{'、'.join(dominant_fonts)}")
        if dominant_font_sizes:
            summary_parts.append(f"- 高频字号：{'、'.join(dominant_font_sizes)}")
        if dominant_colors:
            summary_parts.append(f"- 高频颜色：{'、'.join(dominant_colors)}")
        if reference_image:
            summary_parts.append("- 已提取代表性图片，可作为视觉风格参考")
        summary_parts.extend(common_master_summary.get("summary_lines") or [])
        summary_parts.append("- 采样页摘要：")
        summary_parts.extend([f"  {line}" for line in slide_summaries[:5]])

        return {
            "analysis_summary": "\n".join(summary_parts),
            "reference_image": reference_image,
            "slide_count": slide_count,
            "sampled_slide_count": len(sampled_slides),
            "stable_master_elements": common_master_summary.get("stable_elements") or [],
        }

    def _extract_html_from_response(self, response_content: str) -> str:
        """Extract HTML code from AI response with improved extraction"""
        import re

        logger.info(f"Extracting HTML from response. Content length: {len(response_content)}")

        # Try to extract HTML code block (most common format)
        html_match = re.search(r'```html\s*(.*?)\s*```', response_content, re.DOTALL)
        if html_match:
            extracted = html_match.group(1).strip()
            logger.info(f"Extracted HTML from code block. Length: {len(extracted)}")
            return extracted

        # Try to extract any code block that contains DOCTYPE
        code_block_match = re.search(r'```[a-zA-Z]*\s*(<!DOCTYPE html.*?</html>)\s*```', response_content, re.DOTALL | re.IGNORECASE)
        if code_block_match:
            extracted = code_block_match.group(1).strip()
            logger.info(f"Extracted HTML from generic code block. Length: {len(extracted)}")
            return extracted

        # Try to extract DOCTYPE HTML directly
        doctype_match = re.search(r'<!DOCTYPE html.*?</html>', response_content, re.DOTALL | re.IGNORECASE)
        if doctype_match:
            extracted = doctype_match.group(0).strip()
            logger.info(f"Extracted HTML from direct match. Length: {len(extracted)}")
            return extracted

        # If no specific pattern found, check if the content itself is HTML
        content_stripped = response_content.strip()
        if content_stripped.lower().startswith('<!doctype html') and content_stripped.lower().endswith('</html>'):
            logger.info(f"Content appears to be direct HTML. Length: {len(content_stripped)}")
            return content_stripped

        # Return original content as last resort
        logger.warning(f"Could not extract HTML from response, returning original content. Preview: {response_content[:200]}")
        return response_content.strip()

    def _validate_html_template(self, html_content: str) -> bool:
        """Validate HTML template with improved error reporting"""
        try:
            if not html_content or not html_content.strip():
                logger.error("HTML validation failed: Content is empty")
                return False

            html_lower = html_content.lower().strip()

            # Check basic HTML structure with more flexible validation
            if not html_lower.startswith('<!doctype html'):
                logger.error(f"HTML validation failed: Missing or incorrect DOCTYPE. Content starts with: {html_content[:100]}")
                return False

            if '</html>' not in html_lower:
                logger.error("HTML validation failed: Missing closing </html> tag")
                return False

            # Check required elements with better error reporting
            required_elements = {
                '<head>': '<head',
                '<body>': '<body',
                '<title>': '<title'
            }
            missing_elements = []

            for element_name, element_pattern in required_elements.items():
                if element_pattern not in html_lower:
                    missing_elements.append(element_name)

            if missing_elements:
                logger.error(f"HTML validation failed: Missing required elements: {missing_elements}")
                return False

            logger.info("HTML template validation passed successfully")
            return True

        except Exception as e:
            logger.error(f"HTML validation failed with exception: {e}")
            return False

    async def _generate_preview_image(self, html_template: str) -> str:
        """Generate preview image for template (placeholder implementation)"""
        # This is a placeholder implementation
        placeholder_svg = """
        <svg width="320" height="180" xmlns="http://www.w3.org/2000/svg">
            <rect width="320" height="180" fill="#f3f4f6"/>
            <text x="160" y="90" text-anchor="middle" font-family="Arial" font-size="14" fill="#6b7280">
                模板预览
            </text>
        </svg>
        """
        return f"data:image/svg+xml;base64,{base64.b64encode(placeholder_svg.encode()).decode()}"

    def _extract_style_config(self, html_content: str) -> Dict[str, Any]:
        """Extract style configuration from HTML"""
        import re

        style_config = {
            "dimensions": "1280x720",
            "aspect_ratio": "16:9",
            "framework": "HTML + CSS"
        }

        try:
            # Extract color configuration
            color_matches = re.findall(r'(?:background|color)[^:]*:\s*([^;]+)', html_content, re.IGNORECASE)
            if color_matches:
                style_config["colors"] = list(set(color_matches[:10]))  # Limit to 10 colors

            # Extract font configuration
            font_matches = re.findall(r'font-family[^:]*:\s*([^;]+)', html_content, re.IGNORECASE)
            if font_matches:
                style_config["fonts"] = list(set(font_matches[:5]))  # Limit to 5 fonts

            # Check for frameworks
            if 'tailwind' in html_content.lower():
                style_config["framework"] = "Tailwind CSS"
            elif 'bootstrap' in html_content.lower():
                style_config["framework"] = "Bootstrap"

        except Exception as e:
            logger.warning(f"Failed to extract style config: {e}")

        return style_config

    async def get_templates_by_tags(self, tags: List[str], active_only: bool = True) -> List[Dict[str, Any]]:
        """Get global master templates by tags"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                templates = await db_service.get_global_master_templates_by_tags(
                    tags,
                    active_only,
                    user_id=self.user_id,
                )

                return [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

        except Exception as e:
            logger.error(f"Failed to get global master templates by tags: {e}")
            raise

    async def get_templates_by_tags_paginated(
        self,
        tags: List[str],
        active_only: bool = True,
        page: int = 1,
        page_size: int = 6,
        search: Optional[str] = None
    ) -> Dict[str, Any]:
        """Get global master templates by tags with pagination"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)

                # Calculate offset
                offset = (page - 1) * page_size

                # Get templates with pagination
                templates, total_count = await db_service.get_global_master_templates_by_tags_paginated(
                    tags=tags,
                    active_only=active_only,
                    offset=offset,
                    limit=page_size,
                    search=search,
                    user_id=self.user_id,
                )

                # Calculate pagination info
                total_pages = (total_count + page_size - 1) // page_size
                has_next = page < total_pages
                has_prev = page > 1

                template_list = [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

                return {
                    "templates": template_list,
                    "pagination": {
                        "current_page": page,
                        "page_size": page_size,
                        "total_count": total_count,
                        "total_pages": total_pages,
                        "has_next": has_next,
                        "has_prev": has_prev
                    }
                }
        except Exception as e:
            logger.error(f"Failed to get paginated templates by tags: {e}")
            raise

    async def increment_template_usage(self, template_id: int) -> bool:
        """Increment template usage count"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                return await db_service.increment_global_master_template_usage(
                    template_id,
                    user_id=self.user_id,
                )

        except Exception as e:
            logger.error(f"Failed to increment template usage {template_id}: {e}")
            raise
